"""
Google Drive Reader
Uses googleapiclient directly to read Shared Drive, bypassing LlamaIndex GoogleDriveReader
due to Shared Drive support issues.
"""
import os
import io
import json
import tempfile
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from sync.models import SourceDocument

_EXPORT_MIME = {
    "application/vnd.google-apps.document": (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".docx",
    ),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    ),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    ),
}

_EDITOR_URL = {
    "application/vnd.google-apps.document": "https://docs.google.com/document/d/{id}/edit",
    "application/vnd.google-apps.spreadsheet": "https://docs.google.com/spreadsheets/d/{id}/edit",
    "application/vnd.google-apps.presentation": "https://docs.google.com/presentation/d/{id}/edit",
}


def _build_service():
    token_path = os.environ.get("GOOGLE_DRIVE_TOKEN_PATH", "credentials/token.json")
    with open(token_path) as f:
        info = json.load(f)
    creds = Credentials.from_authorized_user_info(info)
    return build("drive", "v3", credentials=creds)


def _resolve_shortcut(service, file: dict) -> dict | None:
    """Resolve shortcut to target file metadata, keeping the shortcut's name."""
    target_id = file.get("shortcutDetails", {}).get("targetId")
    target_mime = file.get("shortcutDetails", {}).get("targetMimeType")
    if not target_id:
        return None
    return {
        "id": target_id,
        "name": file["name"],
        "mimeType": target_mime,
        "modifiedTime": file.get("modifiedTime", ""),
    }


def _list_files(service, folder_id: str, _seen_ids: set | None = None) -> list[dict]:
    """Recursively list all files in a folder (including subfolders), auto-following shortcuts."""
    if _seen_ids is None:
        _seen_ids = set()

    files = []
    page_token = None
    while True:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            fields="nextPageToken, files(id, name, mimeType, modifiedTime, shortcutDetails)",
            pageToken=page_token,
        ).execute()
        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                files.extend(_list_files(service, f["id"], _seen_ids))
            elif f["mimeType"] == "application/vnd.google-apps.shortcut":
                target = _resolve_shortcut(service, f)
                if target and target["id"] not in _seen_ids:
                    _seen_ids.add(target["id"])
                    if target["mimeType"] == "application/vnd.google-apps.folder":
                        files.extend(_list_files(service, target["id"], _seen_ids))
                    else:
                        files.append(target)
            else:
                if f["id"] not in _seen_ids:
                    _seen_ids.add(f["id"])
                    files.append(f)
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return files


def _read_file_text(service, file: dict) -> str | None:
    """Download and extract text content"""
    mime = file["mimeType"]
    fid = file["id"]

    if mime in _EXPORT_MIME:
        export_mime, ext = _EXPORT_MIME[mime]
        req = service.files().export_media(fileId=fid, mimeType=export_mime)
    elif mime == "application/pdf" or mime.startswith("text/"):
        req = service.files().get_media(fileId=fid, supportsAllDrives=True)
    else:
        return None

    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)

    if mime in _EXPORT_MIME:
        _, ext = _EXPORT_MIME[mime]
        suffix = ext
    elif mime == "application/pdf":
        suffix = ".pdf"
    else:
        suffix = ".txt"

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(buf.read())
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            import pymupdf4llm
            return pymupdf4llm.to_markdown(tmp_path, show_progress=False)
        elif suffix == ".docx":
            import docx
            doc = docx.Document(tmp_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif suffix == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(tmp_path, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
            return "\n".join(rows)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(tmp_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n\n".join(texts)
        else:
            return buf.getvalue().decode("utf-8", errors="ignore")
    finally:
        os.unlink(tmp_path)


def list_google_drive_files() -> tuple:
    """Returns (service, files) — lightweight listing with modifiedTime for Layer 1 check"""
    folder_id = os.environ["GOOGLE_DRIVE_FOLDER_ID"]
    service = _build_service()
    files = _list_files(service, folder_id)
    print(f"  [google_drive] found {len(files)} files")
    return service, files


def fetch_google_drive_file(service, f: dict) -> SourceDocument | None:
    """Download and build SourceDocument for a single file"""
    try:
        text = _read_file_text(service, f)
        if not text or not text.strip():
            return None
        fid = f["id"]
        mime = f["mimeType"]
        url_template = _EDITOR_URL.get(mime, "https://drive.google.com/file/d/{id}/view")
        return SourceDocument(
            text=text,
            metadata={
                "source_type": "google_drive",
                "source_id": f"drive_{fid}",
                "file_name": f["name"],
                "mime_type": mime,
                "modified_time": f.get("modifiedTime", ""),
                "source_updated_at": f.get("modifiedTime", ""),
                "source_url": url_template.format(id=fid),
            },
        )
    except Exception as e:
        print(f"  [google_drive] skip {f['name']}: {e}")
        return None


def fetch_google_drive_documents() -> list[SourceDocument]:
    """Legacy full-fetch function kept for backward compatibility"""
    service, files = list_google_drive_files()
    docs = []
    for f in files:
        doc = fetch_google_drive_file(service, f)
        if doc:
            docs.append(doc)
    return docs